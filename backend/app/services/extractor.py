from pathlib import Path
from pypdf import PdfReader
from docx import Document
from pptx import Presentation


def extract_text(path: str) -> str:
    file_path = Path(path)
    suffix = file_path.suffix.lower()
    if suffix == ".pdf":
        reader = PdfReader(str(file_path))
        return "\n\n".join((page.extract_text() or "") for page in reader.pages).strip()
    if suffix == ".docx":
        doc = Document(str(file_path))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip()).strip()
    if suffix == ".pptx":
        prs = Presentation(str(file_path))
        lines: list[str] = []
        for index, slide in enumerate(prs.slides, start=1):
            lines.append(f"[SLIDE {index}]")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    lines.append(shape.text.strip())
        return "\n".join(lines).strip()
    if suffix in {".txt", ".md"}:
        return file_path.read_text(encoding="utf-8", errors="ignore").strip()
    raise ValueError("Formato não suportado. Use PDF, DOCX, PPTX, TXT ou MD.")
